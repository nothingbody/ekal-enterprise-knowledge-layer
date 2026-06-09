"""File analyzer tool — reads and analyzes uploaded files in-chat.

Supports text extraction from PDF, Word, Excel, CSV, TXT, Markdown, etc.
Files are processed in memory without being added to any knowledge base.
"""

from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any

from app.core.tools.base import BaseTool, ToolResult

logger = logging.getLogger(__name__)

MAX_EXTRACT_CHARS = 15000


class FileAnalyzerTool(BaseTool):
    name = "file_analyzer"
    description = (
        "分析用户上传的文件内容。"
        "可以读取 PDF、Word、Excel、CSV、TXT、Markdown 等格式的文件，"
        "提取文本内容供分析、摘要、翻译、数据提取等用途。"
        "文件不会被存入知识库，仅在当前对话中使用。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "服务器上的文件路径（由系统提供）",
            },
            "task": {
                "type": "string",
                "description": "分析任务：summary（摘要）、extract（提取要点）、translate（翻译）、analyze（通用分析）、raw（返回原文）",
                "enum": ["summary", "extract", "translate", "analyze", "raw"],
                "default": "analyze",
            },
        },
        "required": ["file_path"],
    }

    async def execute(self, **kwargs) -> ToolResult:
        file_path = kwargs.get("file_path", "").strip()
        task = kwargs.get("task", "analyze")

        if not file_path:
            return ToolResult(success=False, error="文件路径不能为空")

        path = Path(file_path)
        if not path.is_file():
            return ToolResult(success=False, error=f"文件不存在: {file_path}")

        try:
            text = await _extract_text(path)
            if not text:
                return ToolResult(success=False, error="无法从文件中提取文本内容")

            text = text[:MAX_EXTRACT_CHARS]
            suffix = path.suffix.lower()
            file_info = f"文件: {path.name} ({suffix}, {_human_size(path.stat().st_size)})"

            if task == "raw":
                return ToolResult(
                    success=True,
                    data={"text": f"{file_info}\n\n{text}"},
                    display_type="text",
                )

            return ToolResult(
                success=True,
                data={
                    "text": f"{file_info}\n\n--- 文件内容 ---\n{text}",
                    "file_name": path.name,
                    "file_size": path.stat().st_size,
                    "task": task,
                },
                display_type="text",
            )
        except Exception as exc:
            logger.warning("File analysis failed for %s: %s", file_path, exc)
            return ToolResult(success=False, error=f"文件分析失败: {str(exc)}")


async def _extract_text(path: Path) -> str:
    """Extract text content from various file formats."""
    import asyncio
    suffix = path.suffix.lower()

    def _sync_extract():
        if suffix in (".txt", ".md", ".markdown", ".log", ".json", ".xml", ".yaml", ".yml", ".csv"):
            return path.read_text(encoding="utf-8", errors="replace")

        if suffix == ".pdf":
            return _extract_pdf(path)

        if suffix == ".docx":
            return _extract_docx(path)

        if suffix == ".doc":
            return _extract_doc_legacy(path)

        if suffix == ".xlsx":
            return _extract_excel(path)

        if suffix == ".xls":
            return _extract_xls_legacy(path)

        if suffix in (".pptx",):
            return _extract_pptx(path)

        if suffix == ".html":
            return _extract_html(path)

        if suffix in (".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs", ".rb", ".php", ".sh", ".sql"):
            return path.read_text(encoding="utf-8", errors="replace")

        return path.read_text(encoding="utf-8", errors="replace")

    return await asyncio.to_thread(_sync_extract)


def _extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages[:100]:
                text = page.extract_text()
                if text:
                    texts.append(text)
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        texts.append(" | ".join(str(cell or "") for cell in row))
        return "\n\n".join(texts)
    except ImportError:
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(str(path))
            return "\n\n".join(page.extract_text() or "" for page in reader.pages[:100])
        except ImportError:
            return ""


def _extract_docx(path: Path) -> str:
    try:
        import docx
        doc = docx.Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return ""


def _extract_doc_legacy(path: Path) -> str:
    """Extract text from legacy .doc format using textract or antiword fallback."""
    try:
        import subprocess
        result = subprocess.run(
            ["antiword", str(path)], capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return f"[无法解析 .doc 格式文件 '{path.name}'，请转换为 .docx 后重试]"


def _extract_excel(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames[:10]:
            sheet = wb[sheet_name]
            parts.append(f"=== 工作表: {sheet_name} ===")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= 200:
                    parts.append(f"... (仅显示前 200 行)")
                    break
                parts.append(" | ".join(str(cell or "") for cell in row))
                row_count += 1
        wb.close()
        return "\n".join(parts)
    except ImportError:
        return ""


def _extract_xls_legacy(path: Path) -> str:
    """Extract text from legacy .xls format using xlrd."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts = []
        for sheet in wb.sheets()[:10]:
            parts.append(f"=== 工作表: {sheet.name} ===")
            for row_idx in range(min(sheet.nrows, 200)):
                row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                parts.append(" | ".join(row))
            if sheet.nrows > 200:
                parts.append("... (仅显示前 200 行)")
        return "\n".join(parts)
    except ImportError:
        return f"[无法解析 .xls 格式文件 '{path.name}'，请转换为 .xlsx 后重试]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                parts.append(f"--- 第 {i} 页 ---\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except ImportError:
        return ""


def _extract_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
        html = path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        return path.read_text(encoding="utf-8", errors="replace")


def _human_size(size: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if size < 1024:
            return f"{size:.1f}{unit}"
        size /= 1024
    return f"{size:.1f}TB"
