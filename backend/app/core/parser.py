import os
from pathlib import Path
from typing import Optional


def _detect_encoding(raw: bytes) -> str:
    """Detect text encoding. Tries chardet first, then falls back to heuristics."""
    # Try chardet / cchardet if available
    try:
        import chardet
        result = chardet.detect(raw[:32768])
        if result and result.get("encoding") and result.get("confidence", 0) > 0.5:
            return result["encoding"]
    except ImportError:
        pass

    # BOM detection
    if raw[:3] == b'\xef\xbb\xbf':
        return "utf-8-sig"
    if raw[:2] in (b'\xff\xfe', b'\xfe\xff'):
        return "utf-16"

    # Try UTF-8 first (strict)
    try:
        raw[:8192].decode("utf-8", errors="strict")
        return "utf-8"
    except UnicodeDecodeError:
        pass

    # Try GB18030 (superset of GBK/GB2312, common for Chinese)
    try:
        raw[:8192].decode("gb18030", errors="strict")
        return "gb18030"
    except UnicodeDecodeError:
        pass

    return "utf-8"


def parse_txt(file_path: str) -> str:
    with open(file_path, "rb") as f:
        raw = f.read()
    enc = _detect_encoding(raw)
    return raw.decode(enc, errors="ignore")


def parse_markdown(file_path: str) -> str:
    return parse_txt(file_path)


def parse_pdf(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    page_count = len(doc)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
        tables = page.find_tables()
        if tables and tables.tables:
            for table in tables.tables:
                try:
                    df = table.to_pandas()
                    text_parts.append(df.to_markdown(index=False))
                except Exception:
                    pass
    doc.close()

    full_text = "\n".join(text_parts)
    # Detect scanned / image-only PDFs: has pages but negligible extracted text.
    # Return empty string so the caller can fall back to LLM vision OCR.
    stripped_len = len(full_text.strip())
    if page_count > 0 and stripped_len < 50:
        return ""
    return full_text


def _doc_chinese_ratio(text: str) -> float:
    """Return the ratio of CJK characters in *text* (0‒1)."""
    import re
    if not text:
        return 0.0
    cjk = len(re.findall(r'[\u4e00-\u9fff]', text))
    return cjk / len(text) if len(text) > 0 else 0.0


def _clean_doc_text(text: str) -> str:
    """Remove Word/WPS binary formatting artefacts from decoded .doc text."""
    import re

    # Strip null bytes (leftover from UTF-16LE → single-byte decode)
    text = text.replace('\x00', '')

    # Word/WPS run‐property tokens  (CJo(aJ>*, CJOJPJQJ^Jo(, mH sH nHtH …)
    text = re.sub(
        r'(?:CJ|PJ|QJ|OJ)(?:OJPJQJ)?(?:\^J)?o?\([^)]{0,30}\)', '', text)
    text = re.sub(r'\b(?:mH|sH|nHtH|nH|tH)\s*', '', text)

    # Field codes: $If, $$If:VTT44l44l06!…
    text = re.sub(r'\$\$?If[^\n]{0,120}', '', text)

    # WPS layout tokens: WD`, VD`, UD], a$$, G$H$, dha$$, …
    text = re.sub(r'(?:WD|VD|UD|WD\?0)[`\'\]\[]\d*', '', text)
    text = re.sub(r'(?:dh)?a\$\$[^\n]{0,60}', '', text)
    text = re.sub(r'G\$H?\$', '', text)

    # Backslash RTF‐like directives
    text = re.sub(r'\\[a-zA-Z]{2,}[\d*]{0,5}', '', text)

    # Keep only meaningful characters (ASCII printable + CJK + common punct)
    text = re.sub(
        r'[^\x20-\x7E\u4e00-\u9fff\u3000-\u303f\uff00-\uffef'
        r'\u00c0-\u024f\n\r\t]', '', text)

    # Per‐line quality filter: drop lines dominated by formatting residue
    clean_lines: list[str] = []
    for line in text.split('\n'):
        stripped = line.strip()
        if not stripped:
            clean_lines.append('')
            continue
        ratio = _doc_chinese_ratio(stripped)
        alpha = len([c for c in stripped if c.isalpha()])
        # Keep line if ≥15 % CJK, or ≥40 % alphabetic (English paragraphs), or short
        if ratio >= 0.15 or (alpha / len(stripped) >= 0.4) or len(stripped) < 6:
            clean_lines.append(stripped)

    text = '\n'.join(clean_lines)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {4,}', '  ', text)
    return text.strip()


def parse_doc(file_path: str) -> str:
    """解析旧版 Word (.doc) 文件。

    1. 检测伪装为 .doc 的 .docx（ZIP 文件头）→ 转 parse_docx。
    2. 通过 olefile 读取 OLE WordDocument 流。
    3. 尝试 UTF-16LE / GB18030 / GBK / UTF-8 多种编码解码，
       按中文字符密度评分选取最佳结果。
    4. 对解码结果做格式化代码过滤和逐行质量过滤。
    """
    import re

    # ── 检测 .docx 伪装为 .doc ──
    try:
        with open(file_path, "rb") as f:
            magic = f.read(4)
        if magic[:2] == b'PK':
            return parse_docx(file_path)
    except Exception:
        pass

    # ── 读取原始字节 ──
    raw: bytes = b''
    ole_raw: bytes = b''

    try:
        import olefile
        ole = olefile.OleFileIO(file_path)
        if ole.exists("WordDocument"):
            ole_raw = ole.openstream("WordDocument").read()
        ole.close()
    except ImportError:
        pass
    except Exception:
        pass

    try:
        with open(file_path, "rb") as f:
            raw = f.read()
    except Exception:
        pass

    # ── 多编码尝试，按中文密度评分 ──
    candidates: list[tuple[float, str]] = []
    encodings = ["utf-16-le", "gb18030", "gbk", "utf-8"]

    for source in (ole_raw, raw):
        if not source:
            continue
        for enc in encodings:
            try:
                decoded = source.decode(enc, errors="ignore")
                cleaned = _clean_doc_text(decoded)
                if len(cleaned) < 30:
                    continue
                score = _doc_chinese_ratio(cleaned)
                candidates.append((score, cleaned))
            except Exception:
                continue

    if not candidates:
        return ""

    # 选取中文密度最高的结果
    candidates.sort(key=lambda c: c[0], reverse=True)
    best = candidates[0][1]
    return best if len(best) > 50 else ""


def parse_docx(file_path: str) -> str:
    # Detect OLE (.doc) files disguised as .docx — redirect to parse_doc
    try:
        with open(file_path, "rb") as f:
            magic = f.read(4)
        if magic[:4] == b'\xd0\xcf\x11\xe0':  # OLE Compound Document magic
            return parse_doc(file_path)
    except Exception:
        pass

    from docx import Document
    try:
        doc = Document(file_path)
    except Exception:
        # python-docx failed (e.g. BadZipFile) — try XML fallback, then parse_doc
        parts = _extract_docx_xml_text(file_path)
        if parts:
            return "\n".join(parts)
        return parse_doc(file_path)

    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        if rows:
            header = " | ".join(rows[0])
            sep = " | ".join(["---"] * len(rows[0]))
            body = "\n".join([" | ".join(r) for r in rows[1:]])
            parts.append(f"{header}\n{sep}\n{body}")

    if not parts:
        parts = _extract_docx_xml_text(file_path)

    return "\n".join(parts)


def _extract_docx_xml_text(file_path: str) -> list:
    """Fallback: extract all text from .docx raw XML.

    Catches content in text boxes, content controls (SDT), shapes, etc.
    that python-docx's doc.paragraphs may miss.
    """
    import zipfile
    from xml.etree import ElementTree as ET

    W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    paragraphs: list[str] = []
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            for name in zf.namelist():
                if not name.startswith("word/") or not name.endswith(".xml"):
                    continue
                tree = ET.parse(zf.open(name))
                for p_elem in tree.iter(f"{W_NS}p"):
                    texts = [t.text for t in p_elem.iter(f"{W_NS}t") if t.text]
                    line = "".join(texts).strip()
                    if line:
                        paragraphs.append(line)
    except Exception:
        pass
    return paragraphs


_MAX_TABLE_ROWS = 50_000  # Safety limit for CSV/Excel row count


def parse_csv(file_path: str) -> str:
    import pandas as pd
    for enc in ("utf-8", "utf-8-sig", "gbk", "gb18030", "latin-1"):
        try:
            df = pd.read_csv(file_path, encoding=enc, nrows=_MAX_TABLE_ROWS)
            return df.to_string(index=False)
        except (UnicodeDecodeError, UnicodeError):
            continue
    df = pd.read_csv(file_path, encoding="utf-8", errors="ignore", nrows=_MAX_TABLE_ROWS)
    return df.to_string(index=False)


def parse_excel(file_path: str) -> str:
    import pandas as pd
    dfs = pd.read_excel(file_path, sheet_name=None, nrows=_MAX_TABLE_ROWS)
    parts = []
    for name, df in dfs.items():
        parts.append(f"## Sheet: {name}\n{df.to_string(index=False)}")
    return "\n\n".join(parts)


def _iter_pptx_shapes(shapes):
    """Recursively yield all shapes, including those inside GroupShape."""
    for shape in shapes:
        yield shape
        if shape.shape_type is not None:
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    yield from _iter_pptx_shapes(shape.shapes)
            except Exception:
                pass


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in _iter_pptx_shapes(slide.shapes):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                if rows:
                    header = " | ".join(rows[0])
                    sep = " | ".join(["---"] * len(rows[0]))
                    body_lines = [" | ".join(r) for r in rows[1:]]
                    texts.append(f"{header}\n{sep}\n" + "\n".join(body_lines))
        if texts:
            parts.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(parts)


def parse_html(file_path: str) -> str:
    from bs4 import BeautifulSoup
    with open(file_path, "rb") as f:
        raw = f.read()
    enc = _detect_encoding(raw)
    html = raw.decode(enc, errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_image_placeholder(file_path: str) -> str:
    """Placeholder parser for image files — returns empty to trigger LLM OCR."""
    return ""


# Extensions recognised as images (need LLM OCR)
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"}

PARSER_MAP = {
    ".txt": parse_txt,
    ".md": parse_markdown,
    ".pdf": parse_pdf,
    ".doc": parse_doc,
    ".docx": parse_docx,
    ".csv": parse_csv,
    ".xlsx": parse_excel,
    ".xls": parse_excel,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm": parse_html,
    # Image formats — placeholder returns empty; actual OCR handled by document_service
    **{ext: _parse_image_placeholder for ext in IMAGE_EXTENSIONS},
}


def parse_document(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    parser = PARSER_MAP.get(ext)
    if parser is None:
        raise ValueError(f"不支持的文件格式: {ext}")
    return parser(file_path)
